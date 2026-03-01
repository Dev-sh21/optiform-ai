import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_styled_presentation():
    # Start with a blank presentation
    prs = Presentation()
    
    # Custom Brand Colors (L&T inspired or modern tech)
    PRIMARY_COLOR = RGBColor(0, 51, 102) # Dark Blue
    ACCENT_COLOR = RGBColor(0, 153, 204) # Light Blue
    TEXT_COLOR = RGBColor(51, 51, 51) # Dark Gray
    WHITE = RGBColor(255, 255, 255)
    
    # Check for images
    img_ai_advisor = "/Users/deveshmishra/.gemini/antigravity/brain/a1c8dfda-f9de-4f88-b271-c2c11a33278b/final_ai_advisor.png"
    img_maintenance = "/Users/deveshmishra/.gemini/antigravity/brain/a1c8dfda-f9de-4f88-b271-c2c11a33278b/final_maintenance.png"
    img_dashboard = "/Users/deveshmishra/Documents/creaTech/docs/visuals/final_v1_dashboard.png"
    img_command = "/Users/deveshmishra/.gemini/antigravity/brain/a1c8dfda-f9de-4f88-b271-c2c11a33278b/command_overview_fresh_1772292891445.png"

    def apply_title_style(title_shape):
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    def apply_body_style(body_shape):
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = TEXT_COLOR
            # Try to add bullet styling if it's a list

    def add_header_footer(slide):
        # Decorative top bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PRIMARY_COLOR
        top_bar.line.color.rgb = PRIMARY_COLOR
        
        # Decorative bottom bar
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = ACCENT_COLOR
        bottom_bar.line.color.rgb = ACCENT_COLOR
        
        # Footer text
        txBox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.35), Inches(4), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Team C.O.R.E. | L&T CreaTech 2026"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Background for title slide
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.color.rgb = PRIMARY_COLOR
    
    title = slide.shapes.title
    title.text = "OptiForm AI"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.size = Pt(60)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Intelligent Formwork Optimization\nProblem Statement 4\nTeam C.O.R.E."
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = ACCENT_COLOR
        p.font.size = Pt(24)

    # --- Slide 2: The Problem ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_header_footer(slide)
    
    title = slide.shapes.title
    title.text = "The Problem: Manual Inefficiency"
    apply_title_style(title)
    
    content = slide.placeholders[1]
    content.text = "Current Challenges in Formwork Management:\n"
    content.text += "• Manual calculation of Bill of Quantities (BoQ) is slow & error-prone.\n"
    content.text += "• Overlooking vertical floor repetitions leads to excess inventory.\n"
    content.text += "• Lack of data-driven 'Rent vs. Buy' analytics inflates project costs.\n"
    content.text += "• Un-tracked carbon footprints from poor asset utilization."
    apply_body_style(content)

    # --- Slide 3: Our Solution ---
    slide = prs.slides.add_slide(slide_layout)
    add_header_footer(slide)
    
    title = slide.shapes.title
    title.text = "Our Solution: OptiForm AI"
    apply_title_style(title)
    
    content = slide.placeholders[1]
    content.text = "Transforming manual workflows into an intelligent, automated pipeline:\n"
    content.text += "• Repetition Analytics: ML clustering groups repeating floor designs.\n"
    content.text += "• Precision Kitting: Instant, optimized BoQ assembly kits generated via OR.\n"
    content.text += "• Financial Engine: Dynamic ROI modeling for Rent vs. Buy decisions.\n"
    content.text += "• Sustainability Metrics: Automated carbon tracking for formwork reuse."
    apply_body_style(content)

    # --- Slide 4: Out-of-the-Box Innovation ---
    slide_layout_two_obj = prs.slide_layouts[3] # Two content
    slide = prs.slides.add_slide(slide_layout_two_obj)
    add_header_footer(slide)
    
    title = slide.shapes.title
    title.text = "Out-of-the-Box Innovation"
    apply_title_style(title)
    
    left_content = slide.placeholders[1]
    left_content.text = "Beyond standard optimization:\n\n"
    left_content.text += "1. AI Design Advisor:\nSuggests architectural standardizations to reduce total BoQ.\n\n"
    left_content.text += "2. Predictive Maintenance:\nTracks asset health to predict failures before they happen."
    apply_body_style(left_content)
    
    # Add Image to right placeholder
    right_placeholder = slide.placeholders[2]
    if os.path.exists(img_ai_advisor):
        slide.shapes.add_picture(img_ai_advisor, right_placeholder.left, right_placeholder.top, width=right_placeholder.width)

    # --- Slide 5: The Analytics Pipeline ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_header_footer(slide)
    
    title = slide.shapes.title
    title.text = "The Analytics Pipeline"
    apply_title_style(title)
    
    content = slide.placeholders[1]
    content.text = "A seamless workflow from design data to final procurement:\n\n"
    content.text += "1️⃣ Data Ingress (Floor Plans & BIM Data)\n"
    content.text += "2️⃣ Repetition Clustering (Machine Learning)\n"
    content.text += "3️⃣ Optimization Engine (BoQ Kitting)\n"
    content.text += "4️⃣ Financial & Decision Engine (ROI & Strategy)"
    apply_body_style(content)

    # --- Slide 6: The Prototype Dashboard ---
    slide_layout_blank = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout_blank)
    add_header_footer(slide)
    
    title = slide.shapes.title
    title.text = "Working Prototype: Command Center"
    apply_title_style(title)
    
    img_to_use = img_dashboard if os.path.exists(img_dashboard) else img_command
    if os.path.exists(img_to_use):
        pic = slide.shapes.add_picture(img_to_use, Inches(0.5), Inches(1.8), width=Inches(9))

    # --- Slide 7: Conclusion & Future Scope ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_header_footer(slide)
    
    title = slide.shapes.title
    title.text = "Value Proposition & Future Scope"
    apply_title_style(title)
    
    content = slide.placeholders[1]
    content.text = "Why OptiForm AI for L&T?\n"
    content.text += "• Drastically reduces planning time and excess inventory holding costs.\n"
    content.text += "• Hard data for Rent vs. Buy decisions protects project margins.\n\n"
    content.text += "Future Roadmap:\n"
    content.text += "• Direct SAP ERP and BIM synchronization (Digital Twin ready).\n"
    content.text += "• Real-time IoT tagging for on-site asset tracking."
    apply_body_style(content)

    # Save
    output_path = "/Users/deveshmishra/Documents/creaTech/OptiForm_AI_LnT_Pitch_Styled.pptx"
    prs.save(output_path)
    print(f"Styled presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_styled_presentation()
