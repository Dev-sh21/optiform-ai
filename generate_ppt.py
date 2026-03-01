import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Check for images
    img_ai_advisor = "/Users/deveshmishra/.gemini/antigravity/brain/a1c8dfda-f9de-4f88-b271-c2c11a33278b/final_ai_advisor.png"
    img_maintenance = "/Users/deveshmishra/.gemini/antigravity/brain/a1c8dfda-f9de-4f88-b271-c2c11a33278b/final_maintenance.png"
    img_dashboard = "/Users/deveshmishra/Documents/creaTech/docs/visuals/final_v1_dashboard.png"
    img_command = "/Users/deveshmishra/.gemini/antigravity/brain/a1c8dfda-f9de-4f88-b271-c2c11a33278b/command_overview_fresh_1772292891445.png"

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "OptiForm AI: Intelligent Formwork Optimization"
    subtitle.text = "Team C.O.R.E.\nProblem Statement 4: Automation of Formwork Kitting & BoQ Optimization\nL&T CreaTech 2026"

    # Slide 2: Problem Understanding
    slide_layout = prs.slide_layouts[1] # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Problem"
    content = slide.placeholders[1]
    content.text = "Current Challenges in Formwork Management:\n"
    content.text += "• Manual, slow, and error-prone calculation of Bill of Quantities (BoQ).\n"
    content.text += "• Sub-optimal reuse of formwork due to missed repetition patterns across floors.\n"
    content.text += "• Lack of data-driven 'Rent vs. Buy' decision making, inflating project costs.\n"
    content.text += "• Un-tracked carbon footprints from poor asset utilization."

    # Slide 3: Our Solution
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Our Solution: OptiForm AI"
    content = slide.placeholders[1]
    content.text = "A data-driven platform designed to transform manual BoQ calculations into an intelligent, automated pipeline:\n"
    content.text += "• Repetition Analytics: Automatically clusters vertical repetition patterns across floors.\n"
    content.text += "• Precision Kitting: Uses operations research to generate optimized BoQ assembly kits.\n"
    content.text += "• Financial Strategy: Analyzes rental bids against procurement to decide Rent vs. Buy.\n"
    content.text += "• Sustainability Metrics: Calculates carbon reductions from optimized material reuse."

    # Slide 4: Key Features & Innovation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Out-of-the-Box Innovation"
    content = slide.placeholders[1]
    content.text = "Beyond standard optimization, OptiForm AI includes:\n"
    content.text += "• AI Design Advisor: Analyzes designs and suggests standardizations to further reduce BoQ.\n"
    content.text += "• Predictive Maintenance Engine: Tracks asset health based on reuse history to predict failures.\n"
    content.text += "• Digital Twin Readiness: Ready to export kit configurations to BIM or AR platforms."

    # Add AI advisor image if exists
    if os.path.exists(img_ai_advisor):
        slide.shapes.add_picture(img_ai_advisor, Inches(5), Inches(2.5), width=Inches(4.5))

    # Slide 5: The Analytics Pipeline
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Analytics Pipeline"
    content = slide.placeholders[1]
    content.text = "Seamless workflow from design to procurement:\n"
    content.text += "1. Data Ingress (Floor plans & designs)\n"
    content.text += "2. Repetition Clustering (Machine Learning)\n"
    content.text += "3. Optimization Engine (PuLP / OR-Tools)\n"
    content.text += "4. BoQ Generation (Kitting)\n"
    content.text += "5. Financial & Decision Engine (ROI Engine)"

    # Slide 6: Financial & Sustainability Impact
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Financial & Sustainability Impact"
    content = slide.placeholders[1]
    content.text = "Demonstrable value for L&T projects:\n"
    content.text += "• Reduced Inventory Costs: Eliminating excess through precise repetition clustering.\n"
    content.text += "• Optimized ROI: Ensuring the right mix of renting and buying based on actual usage durations.\n"
    content.text += "• Lower Carbon Footprint: Quantifying the environmental benefit of maximizing plywood and steel reuse cycles."

    # Slide 7: The Prototype Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    title = slide.shapes.title
    title.text = "The Prototype: Command Center"
    if os.path.exists(img_dashboard):
        slide.shapes.add_picture(img_dashboard, Inches(1), Inches(1.5), width=Inches(8))
    elif os.path.exists(img_command):
        slide.shapes.add_picture(img_command, Inches(1), Inches(1.5), width=Inches(8))

    # Slide 8: Conclusion & Future Scope
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion & Future Scope"
    content = slide.placeholders[1]
    content.text = "Why OptiForm AI?\n"
    content.text += "• Readily solves L&T's Problem Statement 4.\n"
    content.text += "• Enhances speed, accuracy, and profitability of formwork management.\n"
    content.text += "\nFuture integration roadmap:\n"
    content.text += "• Direct SAP ERP and BIM synchronization.\n"
    content.text += "• IoT integration for real-time asset tracking on site."

    output_path = "/Users/deveshmishra/Documents/creaTech/OptiForm_AI_LnT_Pitch.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
